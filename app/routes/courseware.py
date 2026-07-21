import json
import os
import base64
from urllib.parse import quote
from flask import Blueprint, render_template, request, jsonify, Response, current_app
from werkzeug.utils import secure_filename
from app.models import Courseware, KnowledgeItem, Quiz
from app import db
from sqlalchemy import or_

courseware_bp = Blueprint('courseware', __name__)

# 允许上传的图片扩展名
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'svg'}

# 内置背景板定义
BACKGROUND_PRESETS = {
    'seagreen': {'name': '青绿色渐变', 'gradient': 'linear-gradient(135deg, #2E8B57 0%, #00897B 100%)', 'text_color': '#FFFFFF'},
    'blue': {'name': '蓝色渐变', 'gradient': 'linear-gradient(135deg, #1565C0 0%, #42A5F5 100%)', 'text_color': '#FFFFFF'},
    'orange': {'name': '橙色渐变', 'gradient': 'linear-gradient(135deg, #E65100 0%, #FFB74D 100%)', 'text_color': '#FFFFFF'},
    'purple': {'name': '紫色渐变', 'gradient': 'linear-gradient(135deg, #6A1B9A 0%, #BA68C8 100%)', 'text_color': '#FFFFFF'},
    'paper': {'name': '浅色纸面', 'gradient': 'linear-gradient(135deg, #FAFAF5 0%, #F5F5DC 100%)', 'text_color': '#2C3E50'},
    'dark': {'name': '深色主题', 'gradient': 'linear-gradient(135deg, #1A1A2E 0%, #16213E 100%)', 'text_color': '#E8EAED'},
}


def _allowed_image_file(filename):
    """检查是否为允许的图片文件"""
    if not filename:
        return False
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    return ext in ALLOWED_IMAGE_EXTENSIONS


def _get_uploads_dir():
    """获取上传目录绝对路径，不存在则创建"""
    uploads_dir = os.path.join(current_app.static_folder, 'uploads')
    if not os.path.exists(uploads_dir):
        os.makedirs(uploads_dir, exist_ok=True)
    return uploads_dir


@courseware_bp.route('/')
def index():
    """课件管理页面"""
    coursewares = Courseware.query.order_by(Courseware.created_at.desc()).all()
    return render_template('courseware.html', coursewares=coursewares)


@courseware_bp.route('/edit/<int:cw_id>')
def edit(cw_id):
    """课件编辑页面"""
    cw = Courseware.query.get_or_404(cw_id)
    return render_template('courseware_edit.html', courseware=cw, background_presets=BACKGROUND_PRESETS)


@courseware_bp.route('/edit/new')
def edit_new():
    """新建课件编辑页面"""
    mode = request.args.get('mode', 'html')
    return render_template('courseware_edit.html', courseware=None, background_presets=BACKGROUND_PRESETS, default_mode=mode)


@courseware_bp.route('/api/list')
def api_list():
    """课件列表API"""
    coursewares = Courseware.query.order_by(Courseware.created_at.desc()).all()
    result = []
    for cw in coursewares:
        content = {}
        try:
            content = json.loads(cw.content_json) if cw.content_json else {}
        except (json.JSONDecodeError, TypeError):
            pass
        sections = content.get('sections', [])
        slides = content.get('slides', [])
        count = len(slides) if cw.mode == 'ppt' else len(sections)
        result.append({
            'id': cw.id,
            'title': cw.title,
            'description': cw.description,
            'mode': cw.mode or 'html',
            'count': count,
            'background': content.get('background', 'seagreen'),
            'created_at': cw.created_at.isoformat() if cw.created_at else None
        })
    return jsonify({'items': result})


@courseware_bp.route('/api/create', methods=['POST'])
def api_create():
    """创建新课件"""
    data = request.get_json()
    title = data.get('title', '未命名课件')
    description = data.get('description', '')
    mode = data.get('mode', 'html')
    content_json = data.get('content_json', '{}')

    if isinstance(content_json, dict):
        if 'background' not in content_json:
            content_json['background'] = 'seagreen'
        if 'custom_bg_url' not in content_json:
            content_json['custom_bg_url'] = None
        if mode == 'ppt':
            if 'slides' not in content_json:
                content_json['slides'] = []
            if 'slideSize' not in content_json:
                content_json['slideSize'] = {'width': 960, 'height': 540, 'ratio': '16:9'}
        else:
            if 'sections' not in content_json:
                content_json['sections'] = []
        content_json = json.dumps(content_json, ensure_ascii=False)

    cw = Courseware(title=title, description=description, mode=mode, content_json=content_json)
    db.session.add(cw)
    db.session.commit()

    return jsonify({'id': cw.id, 'message': '课件已创建'})


@courseware_bp.route('/api/detail/<int:cw_id>')
def api_detail(cw_id):
    """获取课件详情"""
    cw = Courseware.query.get_or_404(cw_id)
    content = {}
    try:
        content = json.loads(cw.content_json) if cw.content_json else {}
    except (json.JSONDecodeError, TypeError):
        pass
    mode = cw.mode or 'html'
    if mode == 'ppt':
        if 'slides' not in content:
            content['slides'] = []
        if 'slideSize' not in content:
            content['slideSize'] = {'width': 960, 'height': 540, 'ratio': '16:9'}
    else:
        if 'sections' not in content:
            content['sections'] = []
    if 'background' not in content:
        content['background'] = 'seagreen'
    if 'custom_bg_url' not in content:
        content['custom_bg_url'] = None
    return jsonify({
        'id': cw.id,
        'title': cw.title,
        'description': cw.description,
        'mode': mode,
        'content_json': content,
        'created_at': cw.created_at.isoformat() if cw.created_at else None
    })


@courseware_bp.route('/api/update/<int:cw_id>', methods=['PUT'])
def api_update(cw_id):
    """更新课件"""
    cw = Courseware.query.get_or_404(cw_id)
    data = request.get_json()
    if 'title' in data:
        cw.title = data['title']
    if 'description' in data:
        cw.description = data['description']
    if 'mode' in data:
        cw.mode = data['mode']
    if 'content_json' in data:
        if isinstance(data['content_json'], str):
            content = json.loads(data['content_json'])
        else:
            content = data['content_json']
            cw.content_json = json.dumps(content, ensure_ascii=False)
            db.session.commit()
            return jsonify({'message': '课件已更新'})
        cw.content_json = json.dumps(content, ensure_ascii=False)
    db.session.commit()
    return jsonify({'message': '课件已更新'})


@courseware_bp.route('/api/delete/<int:cw_id>', methods=['DELETE'])
def api_delete(cw_id):
    """删除课件"""
    cw = Courseware.query.get_or_404(cw_id)
    db.session.delete(cw)
    db.session.commit()
    return jsonify({'message': '课件已删除'})


@courseware_bp.route('/api/preview/<int:cw_id>')
def api_preview(cw_id):
    """预览课件（返回HTML课件页面）"""
    cw = Courseware.query.get_or_404(cw_id)
    content = {}
    try:
        content = json.loads(cw.content_json) if cw.content_json else {}
    except (json.JSONDecodeError, TypeError):
        pass
    
    if cw.mode == 'ppt':
        return api_ppt_preview(cw_id)
    
    return render_template('courseware_preview.html', courseware=cw, content=content,
                           background_presets=BACKGROUND_PRESETS, standalone=False)


def _image_to_base64(image_path):
    """将本地图片转换为base64 data URL"""
    if not os.path.exists(image_path):
        return None
    ext = image_path.rsplit('.', 1)[-1].lower() if '.' in image_path else 'png'
    mime_map = {
        'png': 'image/png', 'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
        'gif': 'image/gif', 'bmp': 'image/bmp', 'webp': 'image/webp', 'svg': 'image/svg+xml'
    }
    mime = mime_map.get(ext, 'image/png')
    try:
        with open(image_path, 'rb') as f:
            data = base64.b64encode(f.read()).decode('utf-8')
        return f'data:{mime};base64,{data}'
    except Exception:
        return None


@courseware_bp.route('/api/export/<int:cw_id>')
def api_export(cw_id):
    """导出课件为HTML文件下载。

    导出的HTML为独立文件，所有CSS/JS内联，图片转base64，可离线打开。
    """
    cw = Courseware.query.get_or_404(cw_id)
    content = {}
    try:
        content = json.loads(cw.content_json) if cw.content_json else {}
    except (json.JSONDecodeError, TypeError):
        pass

    # 处理 sections 中的图片URL，转换为 base64
    sections = content.get('sections', [])
    static_dir = current_app.static_folder
    for sec in sections:
        # 处理 image 类型的元素
        if sec.get('type') == 'image' and sec.get('extra_data'):
            img_url = sec['extra_data'].get('url') if isinstance(sec['extra_data'], dict) else None
            if img_url:
                # 将 /static/uploads/xxx.png 转换为绝对路径
                if img_url.startswith('/static/'):
                    abs_path = os.path.join(static_dir, img_url.replace('/static/', ''))
                    b64 = _image_to_base64(abs_path)
                    if b64:
                        sec['extra_data'] = dict(sec['extra_data'])
                        sec['extra_data']['url'] = b64
                elif img_url.startswith('data:'):
                    pass  # 已经是 base64
        # 处理自定义背景图
    custom_bg_url = content.get('custom_bg_url')
    if custom_bg_url and custom_bg_url.startswith('/static/'):
        abs_path = os.path.join(static_dir, custom_bg_url.replace('/static/', ''))
        b64 = _image_to_base64(abs_path)
        if b64:
            content['custom_bg_url'] = b64

    html = render_template('courseware_preview.html', courseware=cw, content=content,
                           background_presets=BACKGROUND_PRESETS, standalone=True)
    filename = quote(cw.title) + '.html'
    return Response(
        html,
        mimetype='text/html',
        headers={'Content-Disposition': f"attachment; filename*=UTF-8''{filename}"}
    )


def _render_ppt_slide_to_image(slide_data, slide_size):
    """将PPT幻灯片数据渲染为图片"""
    from PIL import Image, ImageDraw, ImageFont
    import io
    
    width = slide_size.get('width', 960)
    height = slide_size.get('height', 540)
    
    img = Image.new('RGB', (width, height), color='#ffffff')
    draw = ImageDraw.Draw(img)
    
    if slide_data.get('background'):
        bg = slide_data['background']
        if bg.get('type') == 'color':
            try:
                img = Image.new('RGB', (width, height), color=bg.get('value', '#ffffff'))
                draw = ImageDraw.Draw(img)
            except:
                pass
        elif bg.get('type') == 'image':
            try:
                bg_url = bg.get('value', '')
                if bg_url.startswith('/static/'):
                    bg_path = os.path.join(current_app.static_folder, bg_url.replace('/static/', ''))
                    bg_img = Image.open(bg_path)
                    bg_img = bg_img.resize((width, height), Image.LANCZOS)
                    img.paste(bg_img, (0, 0))
                    draw = ImageDraw.Draw(img)
            except:
                pass
    
    elements = slide_data.get('elements', [])
    elements.sort(key=lambda e: e.get('zIndex', 1))
    
    for el in elements:
        x = el.get('x', 0)
        y = el.get('y', 0)
        w = el.get('width', 100)
        h = el.get('height', 50)
        rotation = el.get('rotation', 0)
        opacity = el.get('opacity', 1)
        
        el_img = Image.new('RGBA', (w, h), (0, 0, 0, 0))
        el_draw = ImageDraw.Draw(el_img)
        
        el_type = el.get('type', 'text')
        content = el.get('content', {})
        
        if el_type == 'text':
            text = content.get('content', '')
            font_size = content.get('fontSize', 16)
            color = content.get('color', '#333333')
            bold = content.get('bold', False)
            align = content.get('align', 'left')
            
            try:
                font_path = 'C:/Windows/Fonts/msyh.ttc'
                font = ImageFont.truetype(font_path, font_size)
            except:
                font = ImageFont.load_default()
            
            lines = text.split('\n')
            line_height = font_size + 4
            max_line_width = w - 10
            
            for i, line in enumerate(lines):
                text_width = el_draw.textlength(line, font=font)
                if align == 'center':
                    tx = (w - text_width) / 2
                elif align == 'right':
                    tx = w - text_width - 5
                else:
                    tx = 5
                ty = i * line_height + 5
                
                el_draw.text((tx, ty), line, font=font, fill=color)
        
        elif el_type == 'image':
            img_url = content.get('url', '')
            if img_url:
                try:
                    if img_url.startswith('/static/'):
                        img_path = os.path.join(current_app.static_folder, img_url.replace('/static/', ''))
                        src_img = Image.open(img_path)
                        src_img = src_img.resize((w, h), Image.LANCZOS)
                        el_img.paste(src_img, (0, 0))
                    elif img_url.startswith('data:'):
                        import base64
                        data = img_url.split(',')[1]
                        src_img = Image.open(io.BytesIO(base64.b64decode(data)))
                        src_img = src_img.resize((w, h), Image.LANCZOS)
                        el_img.paste(src_img, (0, 0))
                except:
                    el_draw.rectangle([0, 0, w, h], fill='#cccccc')
                    el_draw.text((5, 5), '图片加载失败', fill='#666666')
        
        elif el_type == 'shape':
            shape = content.get('shape', 'rect')
            fill_color = content.get('fillColor', '#2E8B57')
            stroke_color = content.get('strokeColor', '#1B5E20')
            stroke_width = content.get('strokeWidth', 2)
            
            if shape == 'rect':
                el_draw.rectangle([0, 0, w, h], fill=fill_color, outline=stroke_color, width=stroke_width)
            elif shape == 'circle':
                el_draw.ellipse([0, 0, w, h], fill=fill_color, outline=stroke_color, width=stroke_width)
            elif shape == 'line' or shape == 'arrow':
                el_draw.line([0, h/2, w, h/2], fill=stroke_color, width=stroke_width)
                if shape == 'arrow':
                    arrow_size = 10
                    el_draw.polygon([
                        (w - arrow_size, h/2 - arrow_size/2),
                        (w, h/2),
                        (w - arrow_size, h/2 + arrow_size/2)
                    ], fill=stroke_color)
        
        elif el_type == 'graph':
            graph_data = content.get('graphData', {})
            if graph_data:
                try:
                    graph_img = Image.new('RGB', (w, h), '#f8fcf8')
                    graph_draw = ImageDraw.Draw(graph_img)
                    
                    vertices = graph_data.get('vertices', [])
                    edges = graph_data.get('edges', [])
                    
                    vmap = {v.get('id'): v for v in vertices}
                    
                    for e in edges:
                        src = vmap.get(e.get('source'))
                        tgt = vmap.get(e.get('target'))
                        if src and tgt:
                            x1, y1 = src.get('x', 0), src.get('y', 0)
                            x2, y2 = tgt.get('x', 0), tgt.get('y', 0)
                            color = e.get('color', '#444444')
                            graph_draw.line([x1, y1, x2, y2], fill=color, width=2)
                    
                    for v in vertices:
                        x, y = v.get('x', 0), v.get('y', 0)
                        r = v.get('radius', 22)
                        color = v.get('color', '#2E8B57')
                        label = v.get('label', '')
                        
                        graph_draw.ellipse([x-r, y-r, x+r, y+r], fill=color, outline='white', width=2)
                        try:
                            font = ImageFont.truetype('C:/Windows/Fonts/msyh.ttc', 14)
                        except:
                            font = ImageFont.load_default()
                        label_w = graph_draw.textlength(label, font=font)
                        graph_draw.text((x - label_w/2, y - 5), label, font=font, fill='white')
                    
                    el_img = graph_img.convert('RGBA')
                except:
                    el_draw.rectangle([0, 0, w, h], fill='#f8fcf8')
                    el_draw.text((5, 5), '图形加载失败', fill='#999999')
        
        elif el_type == 'knowledge' or el_type == 'quiz':
            title = el.get('title', '')
            text = content.get('content', '')
            
            try:
                title_font = ImageFont.truetype('C:/Windows/Fonts/msyh.ttc', 14)
                text_font = ImageFont.truetype('C:/Windows/Fonts/msyh.ttc', 12)
            except:
                title_font = text_font = ImageFont.load_default()
            
            el_draw.rectangle([0, 0, w, h], fill='#ffffff', outline='#C8E6C9', width=2)
            
            title_height = 28
            el_draw.rectangle([0, 0, w, title_height], fill='#F1F8E9')
            
            if el_type == 'knowledge':
                badge_color = '#00897B'
                badge_text = '知识条目'
            else:
                badge_color = '#FF8C00'
                badge_text = '试题'
            
            el_draw.rectangle([5, 6, 85, 18], fill=badge_color, radius=9)
            el_draw.text((42, 8), badge_text, font=text_font, fill='white', anchor='mm')
            
            title_y = title_height + 5
            el_draw.text((5, title_y), title, font=title_font, fill='#2E8B57')
            
            text_lines = text.split('\n')[:10]
            ty = title_y + 20
            for line in text_lines:
                el_draw.text((5, ty), line, font=text_font, fill='#333333')
                ty += 16
        
        if opacity != 1:
            el_img = el_img.convert('RGBA')
            alpha = int(opacity * 255)
            el_img.putalpha(alpha)
        
        if rotation != 0:
            el_img = el_img.rotate(rotation, expand=True)
            new_w, new_h = el_img.size
            offset_x = (new_w - w) / 2
            offset_y = (new_h - h) / 2
            x -= offset_x
            y -= offset_y
            w, h = new_w, new_h
        
        img.paste(el_img, (int(x), int(y)), el_img)
    
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


@courseware_bp.route('/api/ppt-preview/<int:cw_id>')
def api_ppt_preview(cw_id):
    """PPT预览（返回PPT预览页面）"""
    cw = Courseware.query.get_or_404(cw_id)
    content = {}
    try:
        content = json.loads(cw.content_json) if cw.content_json else {}
    except (json.JSONDecodeError, TypeError):
        pass
    
    slides_data = content.get('slides', [])
    slide_size = content.get('slideSize', {'width': 960, 'height': 540})
    
    slide_images = []
    for slide in slides_data:
        img_buf = _render_ppt_slide_to_image(slide, slide_size)
        img_data = base64.b64encode(img_buf.read()).decode('utf-8')
        slide_images.append(f'data:image/png;base64,{img_data}')
    
    return render_template('courseware_ppt_preview.html', 
                           courseware=cw, 
                           slides=slide_images,
                           current_index=0,
                           total_slides=len(slide_images))


@courseware_bp.route('/api/ppt-export/<int:cw_id>')
def api_ppt_export(cw_id):
    """导出PPT为PPTX文件"""
    from pptx import Presentation
    from pptx.util import Inches
    
    cw = Courseware.query.get_or_404(cw_id)
    content = {}
    try:
        content = json.loads(cw.content_json) if cw.content_json else {}
    except (json.JSONDecodeError, TypeError):
        pass
    
    slides_data = content.get('slides', [])
    slide_size = content.get('slideSize', {'width': 960, 'height': 540})
    
    prs = Presentation()
    prs.slide_width = Inches(slide_size['width'] / 96)
    prs.slide_height = Inches(slide_size['height'] / 96)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_slide_layout)
        
        img_buf = _render_ppt_slide_to_image(slide_data, slide_size)
        
        left = top = Inches(0)
        width = Inches(slide_size['width'] / 96)
        height = Inches(slide_size['height'] / 96)
        
        slide.shapes.add_picture(img_buf, left, top, width=width, height=height)
    
    output_buf = io.BytesIO()
    prs.save(output_buf)
    output_buf.seek(0)
    
    filename = quote(cw.title) + '.pptx'
    return Response(
        output_buf,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={'Content-Disposition': f"attachment; filename*=UTF-8''{filename}"}
    )


@courseware_bp.route('/api/upload-image', methods=['POST'])
def api_upload_image():
    """上传图片，返回图片URL"""
    if 'file' not in request.files:
        return jsonify({'error': '未找到文件'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': '未选择文件'}), 400
    if not _allowed_image_file(file.filename):
        return jsonify({'error': '不支持的图片格式，支持: ' + ', '.join(ALLOWED_IMAGE_EXTENSIONS)}), 400

    # 生成安全的文件名，添加时间戳防冲突
    import time
    original_name = secure_filename(file.filename)
    if not original_name:
        original_name = 'image.png'
    timestamp = int(time.time() * 1000)
    ext = original_name.rsplit('.', 1)[-1] if '.' in original_name else 'png'
    filename = f'cw_{timestamp}_{original_name}'

    uploads_dir = _get_uploads_dir()
    file_path = os.path.join(uploads_dir, filename)
    file.save(file_path)

    # 返回URL（使用 /static/uploads/ 前缀）
    image_url = f'/static/uploads/{filename}'
    return jsonify({
        'url': image_url,
        'filename': filename,
        'message': '上传成功'
    })


@courseware_bp.route('/api/render-graph', methods=['POST'])
def api_render_graph():
    """接收图形JSON，返回渲染后的SVG字符串

    请求体格式：
    {
        "vertices": [{"id":1,"x":100,"y":100,"label":"A","color":"#2E8B57","radius":22}],
        "edges": [{"id":1,"source":1,"target":2,"directed":false,"weight":3,"color":"#444"}],
        "annotations": [{"id":1,"x":150,"y":80,"text":"标注","color":"#333"}],
        "freeDraws": [{"id":1,"points":[{"x":0,"y":0}],"color":"#2E8B57","width":3}],
        "width": 600,
        "height": 400
    }
    """
    data = request.get_json() or {}
    vertices = data.get('vertices', [])
    edges = data.get('edges', [])
    annotations = data.get('annotations', [])
    free_draws = data.get('freeDraws', [])
    width = data.get('width', 600)
    height = data.get('height', 400)

    # 计算边界以适配viewBox
    all_x = [v.get('x', 0) for v in vertices] + [a.get('x', 0) for a in annotations]
    all_y = [v.get('y', 0) for v in vertices] + [a.get('y', 0) for a in annotations]
    for fd in free_draws:
        for p in fd.get('points', []):
            all_x.append(p.get('x', 0))
            all_y.append(p.get('y', 0))
    if all_x and all_y:
        min_x = max(0, min(all_x) - 40)
        min_y = max(0, min(all_y) - 40)
        max_x = max(width, max(all_x) + 40)
        max_y = max(height, max(all_y) + 40)
        vb_w = max_x - min_x
        vb_h = max_y - min_y
    else:
        min_x, min_y, vb_w, vb_h = 0, 0, width, height

    # 生成SVG
    svg_parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{min_x} {min_y} {vb_w} {vb_h}" '
        f'style="width:100%;height:auto;background:#f8fcf8;">'
    ]

    # 顶点ID -> 顶点信息映射
    vmap = {v.get('id'): v for v in vertices}

    # 绘制边
    for e in edges:
        src = vmap.get(e.get('source'))
        tgt = vmap.get(e.get('target'))
        if not src or not tgt:
            continue
        x1, y1 = src.get('x', 0), src.get('y', 0)
        x2, y2 = tgt.get('x', 0), tgt.get('y', 0)
        # 计算端点偏移，使边不穿过顶点圆心
        r = src.get('radius', 22)
        import math
        dx, dy = x2 - x1, y2 - y1
        dist = math.hypot(dx, dy) or 1
        ux, uy = dx / dist, dy / dist
        sx1, sy1 = x1 + ux * r, y1 + uy * r
        sx2, sy2 = x2 - ux * r, y2 - uy * r
        color = e.get('color', '#444')
        if e.get('directed'):
            # 有向边：带箭头
            svg_parts.append(
                f'<line x1="{sx1:.1f}" y1="{sy1:.1f}" x2="{sx2:.1f}" y2="{sy2:.1f}" '
                f'stroke="{color}" stroke-width="2" marker-end="url(#arrowhead)"/>'
            )
        else:
            svg_parts.append(
                f'<line x1="{sx1:.1f}" y1="{sy1:.1f}" x2="{sx2:.1f}" y2="{sy2:.1f}" '
                f'stroke="{color}" stroke-width="2"/>'
            )
        # 权重标签
        if e.get('weight') is not None:
            mx, my = (x1 + x2) / 2, (y1 + y2) / 2 - 8
            svg_parts.append(
                f'<rect x="{mx-12:.1f}" y="{my-10:.1f}" width="24" height="16" rx="3" '
                f'fill="#FFF8E1" stroke="#FFC107" stroke-width="1"/>'
                f'<text x="{mx:.1f}" y="{my+2:.1f}" text-anchor="middle" '
                f'font-size="11" fill="#333">{e["weight"]}</text>'
            )

    # 绘制自由画笔
    for fd in free_draws:
        pts = fd.get('points', [])
        if len(pts) < 2:
            continue
        path_d = 'M ' + ' L '.join(f'{p.get("x",0):.1f} {p.get("y",0):.1f}' for p in pts)
        svg_parts.append(
            f'<path d="{path_d}" fill="none" stroke="{fd.get("color", "#2E8B57")}" '
            f'stroke-width="{fd.get("width", 3)}" stroke-linecap="round" stroke-linejoin="round"/>'
        )

    # 绘制顶点
    for v in vertices:
        x, y = v.get('x', 0), v.get('y', 0)
        r = v.get('radius', 22)
        color = v.get('color', '#2E8B57')
        label = v.get('label', str(v.get('id', '')))
        svg_parts.append(
            f'<circle cx="{x:.1f}" cy="{y:.1f}" r="{r}" fill="{color}" stroke="#fff" stroke-width="2"/>'
            f'<text x="{x:.1f}" y="{y+5:.1f}" text-anchor="middle" font-size="14" '
            f'fill="#fff" font-weight="600">{label}</text>'
        )

    # 绘制标注
    for a in annotations:
        x, y = a.get('x', 0), a.get('y', 0)
        text = a.get('text', '')
        color = a.get('color', '#333')
        svg_parts.append(
            f'<text x="{x:.1f}" y="{y:.1f}" text-anchor="middle" font-size="13" '
            f'fill="{color}" font-weight="500">{text}</text>'
        )

    # 箭头定义
    svg_parts.insert(1, '<defs><marker id="arrowhead" markerWidth="10" markerHeight="7" '
                        'refX="9" refY="3.5" orient="auto">'
                        '<polygon points="0 0, 10 3.5, 0 7" fill="#444"/></marker></defs>')

    svg_parts.append('</svg>')
    svg = ''.join(svg_parts)
    return jsonify({'svg': svg, 'width': vb_w, 'height': vb_h})


@courseware_bp.route('/api/backgrounds')
def api_backgrounds():
    """获取内置背景板列表"""
    return jsonify({'backgrounds': BACKGROUND_PRESETS})


@courseware_bp.route('/api/knowledge-list')
def api_knowledge_list():
    """获取知识库条目列表（用于选择添加）"""
    chapter = request.args.get('chapter', '')
    section = request.args.get('section', '')
    category = request.args.get('category', '')
    search = request.args.get('search', '')
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 20, type=int)

    query = KnowledgeItem.query
    if chapter:
        query = query.filter_by(chapter=chapter)
    if section:
        query = query.filter_by(section=section)
    if category:
        query = query.filter_by(category=category)
    if search:
        keyword = f'%{search}%'
        query = query.filter(or_(KnowledgeItem.title.like(keyword), KnowledgeItem.content.like(keyword)))

    pagination = query.order_by(KnowledgeItem.chapter, KnowledgeItem.section).paginate(
        page=page, per_page=per_page, error_out=False
    )

    items_data = [{
        'id': item.id,
        'title': item.title,
        'content': item.content,
        'content_summary': (item.content[:150] + '...') if item.content and len(item.content) > 150 else (item.content or ''),
        'chapter': item.chapter,
        'section': item.section,
        'category': item.category
    } for item in pagination.items]

    return jsonify({
        'items': items_data,
        'total': pagination.total,
        'page': pagination.page,
        'pages': pagination.pages,
        'has_next': pagination.has_next
    })


@courseware_bp.route('/api/chapters')
def api_chapters():
    """获取知识库章节列表（包含层级信息）"""
    chapters = db.session.query(
        KnowledgeItem.chapter,
        db.func.count(KnowledgeItem.id)
    ).group_by(KnowledgeItem.chapter).all()

    sections = db.session.query(
        KnowledgeItem.chapter,
        KnowledgeItem.section,
        db.func.count(KnowledgeItem.id)
    ).filter(KnowledgeItem.section != '').group_by(
        KnowledgeItem.chapter, KnowledgeItem.section
    ).all()

    chapter_data = []
    for ch, cnt in chapters:
        if ch:
            ch_sections = []
            for sch, sec, scnt in sections:
                if sch == ch:
                    ch_sections.append({'name': sec, 'count': scnt})
            chapter_data.append({'name': ch, 'count': cnt, 'sections': ch_sections})

    return jsonify({'chapters': chapter_data})


@courseware_bp.route('/api/knowledge-categories')
def api_knowledge_categories():
    """获取知识条目分类列表"""
    categories = db.session.query(
        KnowledgeItem.category,
        db.func.count(KnowledgeItem.id)
    ).group_by(KnowledgeItem.category).all()

    category_labels = {
        'definition': '定义',
        'theorem': '定理',
        'algorithm': '算法',
        'example': '示例',
        'advanced': '进阶',
        'review': '复习',
        'detail': '详细讲解'
    }

    result = []
    for cat, cnt in categories:
        if cat:
            result.append({
                'value': cat,
                'label': category_labels.get(cat, cat),
                'count': cnt
            })
    return jsonify({'categories': result})


@courseware_bp.route('/api/quiz-categories')
def api_quiz_categories():
    """获取题库分类信息（题型、难度）"""
    question_types = db.session.query(
        Quiz.question_type,
        db.func.count(Quiz.id)
    ).group_by(Quiz.question_type).all()

    difficulties = db.session.query(
        Quiz.difficulty,
        db.func.count(Quiz.id)
    ).group_by(Quiz.difficulty).all()

    type_labels = {
        'choice': '选择题',
        'fill': '填空题',
        'proof': '证明题',
        'calc': '计算题',
        'true_false': '判断题',
        'short_answer': '简答题'
    }

    difficulty_labels = {
        'easy': '简单',
        'medium': '中等',
        'hard': '困难'
    }

    result = {
        'question_types': [],
        'difficulties': []
    }

    for qt, cnt in question_types:
        if qt:
            result['question_types'].append({
                'value': qt,
                'label': type_labels.get(qt, qt),
                'count': cnt
            })

    for diff, cnt in difficulties:
        if diff:
            result['difficulties'].append({
                'value': diff,
                'label': difficulty_labels.get(diff, diff),
                'count': cnt
            })

    return jsonify(result)


@courseware_bp.route('/api/quiz-list')
def api_quiz_list():
    """获取试题列表（用于选择添加）"""
    chapter = request.args.get('chapter', '')
    question_type = request.args.get('question_type', '')
    difficulty = request.args.get('difficulty', '')
    search = request.args.get('search', '')
    page = request.args.get('page', 1, type=int)
    per_page = request.args.get('per_page', 20, type=int)

    query = Quiz.query
    if chapter:
        query = query.filter_by(chapter=chapter)
    if question_type:
        query = query.filter_by(question_type=question_type)
    if difficulty:
        query = query.filter_by(difficulty=difficulty)
    if search:
        keyword = f'%{search}%'
        query = query.filter(Quiz.question_text.like(keyword))

    pagination = query.order_by(Quiz.chapter, Quiz.question_number).paginate(
        page=page, per_page=per_page, error_out=False
    )

    items_data = [{
        'id': item.id,
        'question_number': item.question_number,
        'question_text': item.question_text,
        'question_type': item.question_type,
        'options': item.options,
        'answer': item.answer,
        'chapter': item.chapter,
        'difficulty': item.difficulty
    } for item in pagination.items]

    return jsonify({
        'items': items_data,
        'total': pagination.total,
        'page': pagination.page,
        'pages': pagination.pages
    })


@courseware_bp.route('/api/quiz-detail/<int:quiz_id>')
def api_quiz_detail(quiz_id):
    """获取单个试题详情"""
    item = Quiz.query.get_or_404(quiz_id)
    return jsonify({
        'id': item.id,
        'question_number': item.question_number,
        'question_text': item.question_text,
        'question_type': item.question_type,
        'options': item.options,
        'answer': item.answer,
        'chapter': item.chapter,
        'difficulty': item.difficulty
    })
